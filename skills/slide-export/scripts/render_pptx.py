"""render_pptx.py — Folien-Payload -> echtes .pptx (Issue #446, Fixrunde PR #488).

Gleiche Ursache und gleiches Gegenmittel wie bei
``skills/word-export/scripts/render_docx.py``: bis zur zweiten Fixrunde endete
der slide-export-Pfad bei einer JSON-Zwischenrepraesentation, das Deck selbst
entstand nur als Prosa-Anweisung an den Agenten. Damit war AC4 ("Foliensatz
laesst sich in PowerPoint oeffnen") strukturell unbelegbar und der
Nachweisversuch landete als test-eigener Renderer in
``tests/test_slide_export_pptx_render.py``.

Repo-Code erzeugt das Deck jetzt deterministisch -- analog zu ``render_tex.py``
im Geschwister-Skill ``latex-export``. ``document-skills:pptx`` bleibt das
deklarierte Folien-Backend des Plugins (Praeflight in SKILL.md) und uebernimmt
Layout-Verfeinerung auf der hier erzeugten Datei.

Keine Fabrikation: Kapitel ohne erkennbare Kernaussage werden mit LEEREM
Folienrumpf gerendert und auf stderr gemeldet -- der Agent fragt nach, statt
eine Kernaussage zu erfinden (SKILL.md).

Oeffentliche API:
  render_pptx(payload, out_path) -> list[str]   (Meldungen)
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

#: Layout-Indizes im Standard-Template von python-pptx.
_TITLE_LAYOUT = 0
_TITLE_AND_CONTENT_LAYOUT = 1
_CONTENT_PLACEHOLDER_IDX = 1

_RAHMEN_TITEL = {"kolloquium": "Kolloquium", "konferenz": "Konferenzvortrag"}

#: Fehlt python-pptx, ist das eine unvollstaendige Installation -- verstaendliche
#: Meldung statt Traceback (AC6).
_MISSING_PPTX_MESSAGE = (
    "FEHLER: Das Python-Paket 'python-pptx' ist nicht installiert - "
    "es wird deshalb kein Foliensatz erzeugt. Nachinstallieren mit "
    "`pip install python-pptx` (bzw. `bash scripts/setup.sh` fuer die "
    "komplette Plugin-Umgebung)."
)

__all__ = ["render_pptx", "main"]


def _add_content_slide(presentation, layout, title: str, body: str):
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[_CONTENT_PLACEHOLDER_IDX].text_frame.text = body
    return slide


def render_pptx(payload: dict, out_path: Path | str) -> list[str]:
    """Rendert die Payload aus `build_slide_deck.py` in ein echtes `.pptx`.

    Eine Folie je Kapitel-Eintrag (AC4: eine Kernaussage pro Folie). Mit
    `rahmen` = kolloquium/konferenz kommen Deckblatt und Agenda davor.
    """
    import pptx

    slides = payload.get("slides") or []
    rahmen = (payload.get("rahmen") or "").strip().lower()

    presentation = pptx.Presentation()
    content_layout = presentation.slide_layouts[_TITLE_AND_CONTENT_LAYOUT]

    if rahmen in _RAHMEN_TITEL:
        cover = presentation.slides.add_slide(presentation.slide_layouts[_TITLE_LAYOUT])
        cover.shapes.title.text = _RAHMEN_TITEL[rahmen]
        # Untertitel-Platzhalter: aus den Kapiteln, nichts Erfundenes.
        cover.placeholders[_CONTENT_PLACEHOLDER_IDX].text_frame.text = (
            f"{len(slides)} Kapitel" if slides else "Noch keine Kapitel"
        )
        agenda = presentation.slides.add_slide(content_layout)
        agenda.shapes.title.text = "Agenda"
        frame = agenda.placeholders[_CONTENT_PLACEHOLDER_IDX].text_frame
        for index, entry in enumerate(slides):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = entry.get("title", "")

    messages: list[str] = []
    for entry in slides:
        _add_content_slide(
            presentation,
            content_layout,
            entry.get("title", ""),
            entry.get("core_statement", ""),
        )
        if not entry.get("core_statement"):
            messages.append(
                f"Kapitel '{entry.get('source', '?')}' hat keine erkennbare Kernaussage - "
                "beim Nutzer nachfragen statt eine zu erfinden."
            )

    out = Path(out_path)
    if str(out.parent):
        out.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(out))
    return messages


def _build_arg_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description=(
            "Rendert die Payload aus build_slide_deck.py in ein echtes .pptx "
            "(Schritt 4 von /academic-research:slides, Issue #446)."
        ),
    )
    parser.add_argument("--payload", required=True, help="JSON-Payload aus build_slide_deck.py")
    parser.add_argument("--output", required=True, help="Zieldatei (.pptx)")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = _build_arg_parser().parse_args(argv)

    try:
        import pptx  # noqa: F401
    except ImportError:
        print(_MISSING_PPTX_MESSAGE, file=sys.stderr)
        return 1

    try:
        payload = json.loads(Path(args.payload).read_text(encoding="utf-8"))
    except OSError as exc:
        print(f"FEHLER: Payload nicht lesbar ({exc}).", file=sys.stderr)
        return 1
    except json.JSONDecodeError as exc:
        print(f"FEHLER: Payload '{args.payload}' ist kein gueltiges JSON ({exc}).", file=sys.stderr)
        return 1

    try:
        messages = render_pptx(payload, args.output)
    except OSError as exc:
        print(f"FEHLER: {exc}", file=sys.stderr)
        return 1

    print(f"Erzeugt: {args.output} ({len(payload.get('slides') or [])} Folien)")
    for message in messages:
        print(message, file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
